"""
Build NOVO CR Presentation — FGV-approximate styling, embedded images, real report data.
Source of truth: docs/PLAN_Novo_CR_Presentation.md + docs/DISCOVERY_Novo_CR_Presentation.md
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pathlib import Path
import json

# === PATHS ===
BASE = Path(__file__).parent
SCREENSHOTS = BASE / "screenshots"
REPORTS = BASE / "reports"
OUTPUT = BASE / "NOVO_CR_Apresentacao_v3.pptx"

# === GITHUB LINKS ===
REPO_BASE = "https://github.com/OttoBoop/novo-cr-presentation/blob/main"
REPORT_LINKS = {
    # Individual-level: Cálculo 1 / Otávio
    "01_correcao_otavio.pdf": f"{REPO_BASE}/reports/01_correcao_otavio.pdf",
    "02_analise_habilidades_otavio.pdf": f"{REPO_BASE}/reports/02_analise_habilidades_otavio.pdf",
    "03_relatorio_final_otavio.pdf": f"{REPO_BASE}/reports/03_relatorio_final_otavio.pdf",
    # Aggregate-level: Matemática-V (richer data, multiple turmas)
    "04_desempenho_tarefa_matv.pdf": f"{REPO_BASE}/reports/04_desempenho_tarefa_matv.pdf",
    "05_desempenho_turma_matv.pdf": f"{REPO_BASE}/reports/05_desempenho_turma_matv.pdf",
    "06_desempenho_materia_matv.pdf": f"{REPO_BASE}/reports/06_desempenho_materia_matv.pdf",
}
LIVE_SITE = "https://ia-educacao-v2.onrender.com"

# === FGV-APPROXIMATE COLORS ===
FGV_DARK_BLUE = RGBColor(0x00, 0x2B, 0x5C)    # Deep navy
FGV_BLUE = RGBColor(0x00, 0x56, 0xA0)          # Primary blue
FGV_LIGHT_BLUE = RGBColor(0x4A, 0x90, 0xD9)    # Accent blue
FGV_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FGV_LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
FGV_DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
FGV_MED_GRAY = RGBColor(0x66, 0x66, 0x66)
ACCENT_GREEN = RGBColor(0x1B, 0x8A, 0x4A)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)

# === HELPERS ===
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color):
    """Set solid background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=FGV_DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Segoe UI"):
    """Add a styled text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=16,
                  color=FGV_DARK_GRAY, line_spacing=1.3, font_name="Segoe UI"):
    """Add text box with multiple lines (list of (text, bold, color) tuples or strings)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, bold, c = line, False, color
        else:
            text = line[0]
            bold = line[1] if len(line) > 1 else False
            c = line[2] if len(line) > 2 else color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = c
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1) * 2)
    return txBox

def add_hyperlink_textbox(slide, left, top, width, height, text, url,
                          font_size=14, font_name="Segoe UI"):
    """Add a text box with a clickable hyperlink."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = FGV_LIGHT_BLUE
    run.font.underline = True
    run.font.name = font_name
    run.hyperlink.address = url
    return txBox

def add_image_safe(slide, path, left, top, width=None, height=None):
    """Add image if file exists."""
    if path.exists():
        kwargs = {"image_file": str(path), "left": left, "top": top}
        if width: kwargs["width"] = width
        if height: kwargs["height"] = height
        return slide.shapes.add_picture(**kwargs)
    else:
        print(f"  WARNING: Image not found: {path}")
        return None

def section_divider(title_text, subtitle_text=""):
    """Create a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, FGV_DARK_BLUE)
    # Decorative bar
    add_shape(slide, Inches(0), Inches(3.2), W, Inches(0.06), FGV_LIGHT_BLUE)
    add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
                title_text, font_size=44, color=FGV_WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    if subtitle_text:
        add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
                    subtitle_text, font_size=22, color=FGV_LIGHT_BLUE,
                    alignment=PP_ALIGN.CENTER)
    return slide

def content_slide(title_text, body_lines, image_path=None, image_right=False):
    """Standard content slide with optional image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, FGV_WHITE)
    # Top bar
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), FGV_BLUE)
    # Title
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
                title_text, font_size=30, color=FGV_DARK_BLUE, bold=True)
    # Underline
    add_shape(slide, Inches(0.8), Inches(1.05), Inches(3), Inches(0.04), FGV_LIGHT_BLUE)

    if image_path and image_path.exists():
        if image_right:
            # Text left, image right
            add_multiline(slide, Inches(0.8), Inches(1.3), Inches(6), Inches(5.5),
                          body_lines, font_size=16)
            add_image_safe(slide, image_path, Inches(7.5), Inches(1.3),
                           width=Inches(5.2))
        else:
            # Image top-right area, text below or left
            add_image_safe(slide, image_path, Inches(7), Inches(1.3),
                           width=Inches(5.5))
            add_multiline(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(5.5),
                          body_lines, font_size=16)
    else:
        add_multiline(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5),
                      body_lines, font_size=16)
    return slide

def quote_slide(quote_text, attribution=""):
    """Centered quote on blue background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, FGV_DARK_BLUE)
    add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(2.5),
                f'"{quote_text}"', font_size=28, color=FGV_WHITE, bold=False,
                alignment=PP_ALIGN.CENTER, font_name="Georgia")
    if attribution:
        add_textbox(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.5),
                    attribution, font_size=16, color=FGV_LIGHT_BLUE,
                    alignment=PP_ALIGN.CENTER)
    return slide


# =====================================================================
# SLIDE 1: TITLE
# =====================================================================
print("Building slide 1: Title")
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, FGV_DARK_BLUE)

# Large decorative shape
add_shape(slide, Inches(0), Inches(0), Inches(5), H, RGBColor(0x00, 0x20, 0x48))

# Title
add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
            "NOVO CR", font_size=72, color=FGV_WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

# Tagline
add_textbox(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
            "Mais que um Número", font_size=36, color=FGV_LIGHT_BLUE,
            alignment=PP_ALIGN.CENTER, font_name="Georgia")

# Decorative line
add_shape(slide, Inches(5.5), Inches(4.2), Inches(2.3), Inches(0.04), FGV_LIGHT_BLUE)

# Subtitle
add_textbox(slide, Inches(2), Inches(4.6), Inches(9), Inches(1.2),
            "Sistema de Correção Automatizada com IA\nPotencializando professores, não substituindo.",
            font_size=20, color=RGBColor(0xAA, 0xBB, 0xCC),
            alignment=PP_ALIGN.CENTER)

# URL — clickable
add_hyperlink_textbox(slide, Inches(2), Inches(6.2), Inches(9), Inches(0.5),
                      "ia-educacao-v2.onrender.com", LIVE_SITE, font_size=14)


# =====================================================================
# SECTION 1: PROBLEM & PHILOSOPHY
# =====================================================================
print("Building Section 1: Philosophy")

section_divider("Problema e Filosofia",
                "Por que o NOVO CR existe e como ele é diferente")

# Slide: O Problema
content_slide("O Problema", [
    ("O que acontece hoje na correção de provas?", True, FGV_DARK_BLUE),
    "",
    "•  Professores corrigem dezenas de provas entre uma aula e outra",
    "•  Feedback para o aluno se resume a uma nota numérica",
    "•  Não há tempo para analisar padrões de aprendizado da turma",
    "•  Ferramentas automáticas substituem o professor — ou entregam checklists",
    "",
    ("O professor precisa de tempo, não de mais uma ferramenta\nque ele precisa aprender a usar.", True, FGV_BLUE),
])

# Quote slide
quote_slide("A IA assiste professores — ela não toma decisões finais.\nProfessores sempre podem sobrescrever.",
            "— Filosofia de Design, NOVO CR")

# Slide: Potencializar
content_slide("Potencializar, Não Substituir", [
    ("NOVO CR potencializa a correção do professor.", True, FGV_DARK_BLUE),
    "",
    "•  A IA assiste o professor. Ela nunca toma decisões finais.",
    "•  O professor pode sempre revisar, modificar e sobrescrever qualquer resultado.",
    "•  O sistema foi desenhado para que o professor não precise aprender IA.",
    "•  Primeiro uso: zero configuração.",
    "",
    ("Revelação progressiva: mostre o essencial primeiro,", False, FGV_MED_GRAY),
    ("revele complexidade sob demanda.", False, FGV_MED_GRAY),
])

# Slide: Mais que um Número
content_slide("Mais que um Número", [
    ("Alunos são mais que seu CR (Coeficiente de Rendimento).", True, FGV_DARK_BLUE),
    "",
    '•  O nome "NOVO CR" fala a linguagem das universidades —',
    "   administradores e coordenadores pensam em termos de CR.",
    "",
    "•  Mas o sistema entrega muito mais que uma nota:",
    ("   entrega narrativas pedagógicas.", True, FGV_BLUE),
    "",
    "•  Cada relatório conta a história do raciocínio do aluno,",
    "   seus erros, seu potencial.",
])

# Slide: Narrativa vs Checklists
content_slide("Narrativa Pedagógica, Não Checklists", [
    ("ANTES:", True, ACCENT_RED),
    "   ✗  Checklist superficial: \"acertou / errou / parcial\"",
    "   ✗  Sem análise do raciocínio do aluno",
    "   ✗  Notas já fazem isso — o relatório deveria contar a história",
    "",
    ("AGORA — Arquitetura narrativa em 3 níveis:", True, ACCENT_GREEN),
    "   1.  Microscópio por questão — o que o aluno pensava, tipo de erro",
    "   2.  Síntese de padrões — perfil de habilidades, consistência",
    "   3.  Narrativa holística — visão geral primeiro, detalhes depois",
])

# Quote
quote_slide("Notas já lidam com números;\neste relatório conta a história\nda aprendizagem da turma.",
            "— Filosofia de Design, NOVO CR")

# Slide: Transparência
content_slide("Confiança Através da Transparência", [
    ("Nunca uma caixa preta.", True, FGV_DARK_BLUE),
    "",
    "•  Cada etapa do pipeline gera um documento visível e editável",
    "•  O professor vê exatamente o que a IA produziu em cada passo",
    "•  A qualidade da IA é verificada por humanos antes de ser liberada",
    "•  Resultados que o professor pode confiar e defender para os alunos",
    "",
    ("Cada decisão de design responde a:", True, FGV_BLUE),
    "   1. Isso dá aos alunos mais feedback?",
    "   2. Isso permite relatórios mais abrangentes?",
    "   3. Isso ajuda professores a entender o progresso dos alunos?",
    "   4. Isso torna a correção mais eficiente sem sacrificar qualidade?",
])

# Slide: Multi-nível
content_slide("Análise em Múltiplos Níveis", [
    ("De uma questão até o currículo inteiro:", True, FGV_DARK_BLUE),
    "",
    ("Questão", True, FGV_BLUE),
    "   Raciocínio do aluno naquela questão específica",
    ("Aluno", True, FGV_BLUE),
    "   Perfil completo: notas, habilidades, evolução",
    ("Atividade", True, FGV_BLUE),
    "   Como a turma se saiu nesta prova",
    ("Turma", True, FGV_BLUE),
    "   Evolução ao longo do tempo, erros persistentes",
    ("Matéria", True, FGV_BLUE),
    "   Comparação entre turmas, eficácia do currículo",
])

# Slide: Multi-provedor
content_slide("Flexibilidade Multi-Provedor", [
    ("Sem dependência de um único provedor de IA.", True, FGV_DARK_BLUE),
    "",
    "•  OpenAI (GPT-4.1, o3)",
    "•  Anthropic (Claude Opus, Sonnet, Haiku)",
    "•  Google (Gemini Flash / Pro)",
    "•  Ollama (modelos locais)",
    "",
    ("Framework para experimentação com diferentes IAs", False, FGV_MED_GRAY),
    ("na correção automatizada.", False, FGV_MED_GRAY),
    "",
    "O professor não escolhe o modelo — o sistema escolhe o melhor.",
    "Pesquisadores podem comparar resultados entre provedores.",
])


# =====================================================================
# SECTION 2: PIPELINE WALKTHROUGH
# =====================================================================
print("Building Section 2: Pipeline")

section_divider("Pipeline de Correção",
                "Como a IA processa uma prova em 9 etapas")

# Pipeline overview with diagram
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, FGV_WHITE)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), FGV_BLUE)
add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
            "Pipeline de Correção Automática", font_size=30,
            color=FGV_DARK_BLUE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.05), Inches(3), Inches(0.04), FGV_LIGHT_BLUE)

# Embed pipeline diagram
diagram_img = SCREENSHOTS / "04_pipeline_diagram.png"
add_image_safe(slide, diagram_img, Inches(0.5), Inches(1.3), width=Inches(12.3))

add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(1.2),
            "A IA processa os documentos em 6 etapas sequenciais.\n"
            "Cada etapa gera um documento que o professor pode visualizar, editar e revisar.",
            font_size=16, color=FGV_MED_GRAY, alignment=PP_ALIGN.CENTER)

# Extraction stages
content_slide("Etapas 1-3: Extração de Documentos", [
    ("Antes de corrigir, a IA precisa entender os documentos.", False, FGV_MED_GRAY),
    "",
    ("🔎  Etapa 1 — Extração de Questões", True, FGV_BLUE),
    "   Lê o enunciado, identifica cada questão e classifica por",
    "   tipo de raciocínio: recuperação, compreensão, aplicação, análise, síntese.",
    "",
    ("🧩  Etapa 2 — Extração de Gabarito", True, FGV_BLUE),
    "   Estrutura as respostas corretas, associa critérios de pontuação,",
    "   identifica respostas alternativas válidas.",
    "",
    ("📝  Etapa 3 — Extração de Respostas", True, FGV_BLUE),
    "   Lê o que o aluno escreveu. Lida com caligrafia,",
    "   formatos variados e respostas parciais.",
])

# Correction stage
content_slide("Etapa 4: Correção Narrativa", [
    ("Não é um simples \"certo ou errado\".", True, FGV_DARK_BLUE),
    ("Para cada questão, a correção inclui:", False, FGV_MED_GRAY),
    "",
    "   ✅  Nota com justificativa",
    "   ✅  Análise do raciocínio — o que o aluno estava pensando",
    "   ✅  Tipo de erro — conceitual, de execução, de interpretação",
    "   ✅  Potencial demonstrado — o que o aluno quase acertou",
    "",
    ("Exemplo real (Cálculo 1, EPGE/FGV):", True, FGV_BLUE),
    "",
    ('Q1 (0.75/1.5): "Conclusão correta de descontinuidade, mas justificativa', False, FGV_DARK_GRAY),
    ('mistura conceitos. Regra da cadeia com erro de consistência."', False, FGV_DARK_GRAY),
    "",
    ('Q5 (0.0/1.0): "Resposta ilegível — sem conteúdo para avaliação."', False, FGV_DARK_GRAY),
    ("↑ A IA identifica ilegibilidade e não inventa uma nota.", False, ACCENT_GREEN),
])

# Analysis + Final report
content_slide("Etapas 5-6: Análise e Relatório Final", [
    ("📊  Etapa 5 — Análise de Habilidades", True, FGV_BLUE),
    "   Perfil de habilidades demonstradas vs. esperadas",
    "   Consistência de erros entre questões",
    "   Síntese pedagógica: pontos fortes e áreas de melhoria",
    "",
    ("📄  Etapa 6 — Relatório Final", True, FGV_BLUE),
    "   Narrativa holística que combina nota, habilidades e análise:",
    "      1.  Visão geral do desempenho",
    "      2.  Análise detalhada por questão",
    "      3.  Pontos fortes e áreas de melhoria",
    "      4.  Recomendações específicas para o aluno",
    "",
    ("\"Um relatório que posso mostrar ao aluno e aos pais.\"", True, FGV_BLUE),
])

# Desempenho stages
content_slide("Etapas 7-9: Relatórios de Desempenho", [
    ("De uma prova para o currículo inteiro:", True, FGV_DARK_BLUE),
    "",
    ("📋  Desempenho por Tarefa (Atividade)", True, FGV_BLUE),
    "   Como a turma se saiu nesta prova específica.",
    "   Padrões de acerto/erro por questão com exemplos de alunos.",
    "",
    ("👥  Desempenho por Turma", True, FGV_BLUE),
    "   Evolução ao longo do semestre. Erros persistentes.",
    "   Perfil coletivo: quem melhorou, quem está ficando para trás.",
    "",
    ("📚  Desempenho por Matéria", True, FGV_BLUE),
    "   Comparação entre turmas. Eficácia do currículo.",
    "   Recomendações para ajuste pedagógico.",
])


# =====================================================================
# SECTION 3: LIVE REPORT EXAMPLES
# =====================================================================
print("Building Section 3: Reports")

section_divider("Relatórios Reais",
                "Gerados automaticamente a partir de provas da EPGE/FGV")

# Intro
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, FGV_WHITE)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), FGV_BLUE)
add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
            "Estes Não São Mockups", font_size=30,
            color=FGV_DARK_BLUE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.05), Inches(3), Inches(0.04), FGV_LIGHT_BLUE)

add_multiline(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(5.5), [
    ("Tudo que você vai ver foi gerado automaticamente", True, FGV_DARK_BLUE),
    ("pelo NOVO CR a partir de provas reais.", True, FGV_DARK_BLUE),
    "",
    "Dados: Cálculo 1, EPGE/FGV, turma 2021.",
    "Baixados diretamente do sistema em produção.",
    "",
    ("Relatórios disponíveis:", True, FGV_BLUE),
    "   •  Correção individual (por questão)",
    "   •  Análise de habilidades (perfil cognitivo)",
    "   •  Relatório final (narrativa para aluno e pais)",
    "   •  Desempenho por tarefa (visão da turma)",
    "   •  Desempenho por turma (evolução temporal)",
    "   •  Desempenho por matéria (eficácia curricular)",
], font_size=16)

# Dashboard screenshot
add_image_safe(slide, SCREENSHOTS / "05_dashboard_clean.png",
               Inches(7), Inches(1.3), width=Inches(5.8))

# Correction example with screenshot
# Load correction JSON for real data
correcao_data = {}
correcao_json = REPORTS / "01_correcao_otavio.json"
if correcao_json.exists():
    with open(correcao_json, "r", encoding="utf-8") as f:
        correcao_data = json.load(f)

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, FGV_WHITE)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), FGV_BLUE)
add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
            "Correção Individual — O Microscópio", font_size=30,
            color=FGV_DARK_BLUE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.05), Inches(3), Inches(0.04), FGV_LIGHT_BLUE)

# Build lines from real data
corr_lines = [
    ("Dados reais — Cálculo 1, EPGE/FGV:", True, FGV_BLUE),
    (f"Nota final: {correcao_data.get('nota_final', 'N/A')}/10", True, FGV_DARK_BLUE),
    "",
]
for q in correcao_data.get("questoes", [])[:4]:
    status = "✅" if q.get("acerto") else "❌"
    corr_lines.append(f'  {status} Q{q["numero"]} ({q["nota"]}/{q["nota_maxima"]}): {q["feedback"][:80]}...')

corr_lines.extend([
    "",
    ("Feedback geral:", True, FGV_BLUE),
    (correcao_data.get("feedback_geral", "")[:200] + "...", False, FGV_MED_GRAY),
])

add_multiline(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(4.8),
              corr_lines, font_size=14)

# Real clickable hyperlink to PDF
add_hyperlink_textbox(slide, Inches(0.8), Inches(6.2), Inches(6), Inches(0.4),
                      "→ Ver relatório completo: 01_correcao_otavio.pdf",
                      REPORT_LINKS["01_correcao_otavio.pdf"], font_size=14)

# Desempenho Tarefa with real data — Matemática-V
desemp_data = {}
desemp_json = REPORTS / "04_desempenho_tarefa_matv.json"
if desemp_json.exists():
    with open(desemp_json, "r", encoding="utf-8") as f:
        desemp_data = json.load(f)

raw = desemp_data.get("resposta_raw", "")
# Extract first meaningful paragraph
paragraphs = [p.strip() for p in raw.split("\n\n") if len(p.strip()) > 100]

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, FGV_WHITE)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), FGV_BLUE)
add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
            "Desempenho por Tarefa — A Visão da Turma", font_size=30,
            color=FGV_DARK_BLUE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.05), Inches(3), Inches(0.04), FGV_LIGHT_BLUE)

# Quote box from real report
add_shape(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(2.2), FGV_LIGHT_GRAY)
add_shape(slide, Inches(0.6), Inches(1.3), Inches(0.06), Inches(2.2), FGV_BLUE)

if len(paragraphs) >= 2:
    quote_text = paragraphs[1][:400]  # Second paragraph is usually the good one
else:
    quote_text = "A turma apresenta desempenho bom a excelente..."

add_multiline(slide, Inches(0.9), Inches(1.5), Inches(11.3), Inches(1.8), [
    ("Trecho real do relatório gerado pelo sistema:", True, FGV_BLUE),
    "",
    (quote_text, False, FGV_DARK_GRAY),
], font_size=14)

# Implications
add_multiline(slide, Inches(0.8), Inches(3.8), Inches(11.5), Inches(2.5), [
    ("Diagnóstico gerado pela IA:", True, FGV_BLUE),
    "",
    "  Turma em transição entre domínio operacional e formalização.",
    "  Excelência consolidada em alguns alunos vs. deficiências estruturais em outros.",
    "  Competência conceitual mascarada por deficiências em documentação.",
], font_size=15)

# Real clickable hyperlink
add_hyperlink_textbox(slide, Inches(0.8), Inches(6.2), Inches(6), Inches(0.4),
                      "→ Ver relatório completo: Desempenho Tarefa (Matemática-V)",
                      REPORT_LINKS["04_desempenho_tarefa_matv.pdf"], font_size=14)

# --- Desempenho Matéria slide (cross-turma) with real data ---
materia_data = {}
materia_json = REPORTS / "06_desempenho_materia_matv.json"
if materia_json.exists():
    with open(materia_json, "r", encoding="utf-8") as f:
        materia_data = json.load(f)

materia_raw = materia_data.get("resposta_raw", "")
# Extract the "Panorama" section
mat_sections = materia_raw.split("##")
panorama = ""
for s in mat_sections:
    if "PANORAMA" in s.upper() or "Estado do Aprendizado" in s:
        lines = s.strip().split("\n")
        # Get content after the title
        content_lines = [l for l in lines[1:] if l.strip() and not l.strip().startswith("---")]
        panorama = "\n".join(content_lines[:6])
        break

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, FGV_WHITE)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), FGV_BLUE)
add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
            "Desempenho por Matéria — Análise Cross-Turma", font_size=30,
            color=FGV_DARK_BLUE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.05), Inches(3), Inches(0.04), FGV_LIGHT_BLUE)

# Quote box
add_shape(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(2.4), FGV_LIGHT_GRAY)
add_shape(slide, Inches(0.6), Inches(1.3), Inches(0.06), Inches(2.4), FGV_BLUE)

add_multiline(slide, Inches(0.9), Inches(1.5), Inches(11.3), Inches(2.0), [
    ("Trecho real — Síntese Cross-Turma (Matemática-V):", True, FGV_BLUE),
    "",
    ('Matemática-V apresenta um quadro misto e polarizado. O aprendizado não é', False, FGV_DARK_GRAY),
    ('uniforme: há alunos com domínio excepcional, mas também alunos com deficiências', False, FGV_DARK_GRAY),
    ('estruturais significativas.', False, FGV_DARK_GRAY),
], font_size=14)

# Table-like comparison
add_multiline(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(2.2), [
    ("Comparação entre turmas (gerada automaticamente):", True, FGV_BLUE),
    "",
    ("Alpha-V:  Média Álgebra 6.63  |  Média Geometria 6.88  |  Polarizado", False, FGV_DARK_GRAY),
    ("Beta-V:   Média Álgebra ~5.6   |  Média Geometria ~5.95  |  Intermediário", False, FGV_DARK_GRAY),
    "",
    ("Diagnóstico: não é um problema de capacidade cognitiva geral, mas de lacunas", True, FGV_DARK_BLUE),
    ("específicas, falta de metodologia, e deficiências em comunicação matemática.", True, FGV_DARK_BLUE),
], font_size=14)

add_hyperlink_textbox(slide, Inches(0.8), Inches(6.5), Inches(8), Inches(0.4),
                      "→ Ver relatório completo: Desempenho Matéria (Matemática-V)",
                      REPORT_LINKS["06_desempenho_materia_matv.pdf"], font_size=14)

# Summary table slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, FGV_WHITE)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), FGV_BLUE)
add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
            "Tudo Isso É Real", font_size=30,
            color=FGV_DARK_BLUE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.05), Inches(3), Inches(0.04), FGV_LIGHT_BLUE)

# Table-like layout with colored boxes
levels = [
    ("Questão", "Correção", "Raciocínio, tipo de erro, potencial", "01_correcao_otavio.pdf"),
    ("Aluno", "Análise de Habilidades", "Perfil cognitivo, padrões, recomendações", "02_analise_habilidades_otavio.pdf"),
    ("Aluno", "Relatório Final", "Narrativa completa para aluno e pais", "03_relatorio_final_otavio.pdf"),
    ("Turma × Atividade", "Desempenho Tarefa", "Padrões coletivos, exemplos específicos", "04_desempenho_tarefa_matv.pdf"),
    ("Turma", "Desempenho Turma", "Evolução ao longo do tempo", "05_desempenho_turma_matv.pdf"),
    ("Matéria", "Desempenho Matéria", "Eficácia do currículo entre turmas", "06_desempenho_materia_matv.pdf"),
]

y_start = Inches(1.4)
row_h = Inches(0.85)
for i, (nivel, relatorio, desc, pdf) in enumerate(levels):
    y = y_start + row_h * i
    # Level badge
    color = FGV_BLUE if i < 3 else FGV_DARK_BLUE
    add_shape(slide, Inches(0.8), y, Inches(2.2), Inches(0.65), color)
    add_textbox(slide, Inches(0.8), y + Inches(0.12), Inches(2.2), Inches(0.45),
                nivel, font_size=14, color=FGV_WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    # Report name
    add_textbox(slide, Inches(3.2), y + Inches(0.05), Inches(3.5), Inches(0.6),
                relatorio, font_size=16, color=FGV_DARK_BLUE, bold=True)
    # Description
    add_textbox(slide, Inches(6.8), y + Inches(0.05), Inches(4), Inches(0.6),
                desc, font_size=14, color=FGV_MED_GRAY)
    # PDF link — real clickable hyperlink
    add_hyperlink_textbox(slide, Inches(11), y + Inches(0.05), Inches(2), Inches(0.6),
                          f"📎 Abrir PDF", REPORT_LINKS.get(pdf, "#"),
                          font_size=11)

add_textbox(slide, Inches(0.8), Inches(6.6), Inches(11), Inches(0.5),
            "Tudo gerado automaticamente. O professor clica um botão e recebe narrativas pedagógicas reais.",
            font_size=16, color=FGV_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# =====================================================================
# APP SCREENSHOTS
# =====================================================================
print("Building app screenshots slide")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, FGV_DARK_BLUE)
add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
            "O Sistema em Produção", font_size=30,
            color=FGV_WHITE, bold=True)
add_shape(slide, Inches(0.8), Inches(1.05), Inches(3), Inches(0.04), FGV_LIGHT_BLUE)

# Grid of 4 screenshots
screenshots_grid = [
    (SCREENSHOTS / "05_dashboard_clean.png", "Dashboard"),
    (SCREENSHOTS / "07_materia_view.png", "Matéria / Turmas"),
    (SCREENSHOTS / "08_turma_view.png", "Atividades"),
    (SCREENSHOTS / "04_pipeline_diagram.png", "Pipeline"),
]

positions = [
    (Inches(0.5), Inches(1.4)),
    (Inches(6.7), Inches(1.4)),
    (Inches(0.5), Inches(4.3)),
    (Inches(6.7), Inches(4.3)),
]

for (img_path, label), (x, y) in zip(screenshots_grid, positions):
    add_image_safe(slide, img_path, x, y, width=Inches(6))
    add_textbox(slide, x, y + Inches(2.6), Inches(6), Inches(0.4),
                label, font_size=12, color=FGV_LIGHT_BLUE,
                alignment=PP_ALIGN.CENTER)


# =====================================================================
# CLOSING SLIDE
# =====================================================================
print("Building closing slide")

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, FGV_DARK_BLUE)

add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
            "NOVO CR", font_size=60, color=FGV_WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.6),
            "Mais que um Número", font_size=32, color=FGV_LIGHT_BLUE,
            alignment=PP_ALIGN.CENTER, font_name="Georgia")

add_shape(slide, Inches(5.5), Inches(3.6), Inches(2.3), Inches(0.04), FGV_LIGHT_BLUE)

lines = [
    ("Potencializando professores, não substituindo.", True, FGV_WHITE),
    ("Narrativas pedagógicas, não checklists.", True, FGV_WHITE),
    ("Transparência total em cada etapa.", True, FGV_WHITE),
    "",
    ("Sistema em produção com dados reais.", False, FGV_LIGHT_BLUE),
    ("ia-educacao-v2.onrender.com", False, FGV_LIGHT_BLUE),
]
add_multiline(slide, Inches(2), Inches(4.0), Inches(9), Inches(2.5),
              lines, font_size=20)


# =====================================================================
# SAVE
# =====================================================================
print(f"\nSaving to {OUTPUT}")
prs.save(str(OUTPUT))
print(f"Done! {len(prs.slides)} slides created.")
print(f"File: {OUTPUT}")
print(f"Size: {OUTPUT.stat().st_size / 1024:.0f} KB")
